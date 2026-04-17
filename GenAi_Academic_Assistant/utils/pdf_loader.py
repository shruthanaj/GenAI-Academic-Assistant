"""Document loader utilities - supports PDF, DOCX, XLSX, CSV, PPTX, Markdown, JSON, TXT."""
from __future__ import annotations
import io
import json
import csv
from pathlib import Path

try:
    import PyPDF2
    _HAS_PYPDF = True
except ImportError:
    _HAS_PYPDF = False

try:
    from docx import Document
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False

try:
    from openpyxl import load_workbook
    _HAS_XLSX = True
except ImportError:
    _HAS_XLSX = False

try:
    from pptx import Presentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False


def load_pdf_text(file_obj) -> str:
    """Extract all text from a PDF file-like object."""
    if not _HAS_PYPDF:
        return "[PyPDF2 not installed. Run: pip install PyPDF2]"
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(file_obj.read()))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages)
    except Exception as exc:
        return f"[PDF read error: {exc}]"


def load_docx_text(file_obj) -> str:
    """Extract text from DOCX file."""
    if not _HAS_DOCX:
        return "[python-docx not installed. Run: pip install python-docx]"
    try:
        doc = Document(io.BytesIO(file_obj.read()))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        # Also extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)
        return "\n\n".join(paragraphs)
    except Exception as exc:
        return f"[DOCX read error: {exc}]"


def load_xlsx_text(file_obj) -> str:
    """Extract text from Excel file (XLSX)."""
    if not _HAS_XLSX:
        return "[openpyxl not installed. Run: pip install openpyxl]"
    try:
        workbook = load_workbook(io.BytesIO(file_obj.read()), data_only=True)
        all_text = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            all_text.append(f"\n=== Sheet: {sheet_name} ===\n")
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    all_text.append(row_text)
        return "\n".join(all_text)
    except Exception as exc:
        return f"[XLSX read error: {exc}]"


def load_csv_text(file_obj) -> str:
    """Extract text from CSV file."""
    try:
        content = file_obj.read().decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(content))
        rows = [" | ".join(row) for row in reader if any(row)]
        return "\n".join(rows)
    except Exception as exc:
        return f"[CSV read error: {exc}]"


def load_pptx_text(file_obj) -> str:
    """Extract text from PowerPoint file (PPTX)."""
    if not _HAS_PPTX:
        return "[python-pptx not installed. Run: pip install python-pptx]"
    try:
        presentation = Presentation(io.BytesIO(file_obj.read()))
        all_text = []
        for i, slide in enumerate(presentation.slides, 1):
            all_text.append(f"\n=== Slide {i} ===\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    all_text.append(shape.text)
        return "\n".join(all_text)
    except Exception as exc:
        return f"[PPTX read error: {exc}]"


def load_json_text(file_obj) -> str:
    """Extract text from JSON file."""
    try:
        content = file_obj.read().decode("utf-8", errors="ignore")
        data = json.loads(content)
        # Convert JSON to readable text
        return json.dumps(data, indent=2)
    except Exception as exc:
        return f"[JSON read error: {exc}]"


def load_markdown_text(file_obj) -> str:
    """Load Markdown file as-is."""
    try:
        return file_obj.read().decode("utf-8", errors="ignore")
    except Exception as exc:
        return f"[Markdown read error: {exc}]"


def load_text_file(file_obj) -> str:
    """Load plain text file."""
    try:
        return file_obj.read().decode("utf-8", errors="ignore")
    except Exception as exc:
        return f"[Text read error: {exc}]"


def load_file_text(file_obj, file_name: str) -> str:
    """
    Auto-detect file type and extract text.

    Supported formats: PDF, DOCX, XLSX, CSV, PPTX, JSON, Markdown, TXT
    """
    file_ext = Path(file_name).suffix.lower()

    # Reset file pointer to beginning
    file_obj.seek(0)

    if file_ext == ".pdf":
        return load_pdf_text(file_obj)
    elif file_ext == ".docx":
        return load_docx_text(file_obj)
    elif file_ext == ".xlsx":
        return load_xlsx_text(file_obj)
    elif file_ext == ".csv":
        return load_csv_text(file_obj)
    elif file_ext == ".pptx":
        return load_pptx_text(file_obj)
    elif file_ext == ".json":
        return load_json_text(file_obj)
    elif file_ext in [".md", ".markdown"]:
        return load_markdown_text(file_obj)
    elif file_ext == ".txt":
        return load_text_file(file_obj)
    else:
        # Try as plain text by default
        return load_text_file(file_obj)
